"""Rebuild the AIGRE architecture slide as native, editable PowerPoint shapes
(rectangles, connectors, text boxes) instead of a picture -- so it opens in
Google Slides / PowerPoint / Keynote with everything selectable and movable,
matching the layout of the SVG artifact."""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

# ---- coordinate system: same 1600x940 units as the SVG viewBox -----------
SCALE = 7620  # EMU per unit -> 1600 units = 13.333in, 940 units = 7.833in

def E(v):
    return Emu(int(round(v * SCALE)))

# ---- palette (same tokens as the SVG) -------------------------------------
PAPER = RGBColor(0xF5, 0xF1, 0xE7)
PAPER_LINE = RGBColor(0xE4, 0xDC, 0xC8)
NAVY_900 = RGBColor(0x14, 0x24, 0x3A)
NAVY_700 = RGBColor(0x1B, 0x3A, 0x57)
NAVY_500 = RGBColor(0x3F, 0x5A, 0x78)
NAVY_300 = RGBColor(0x7C, 0x90, 0xA5)
AMBER_600 = RGBColor(0xC1, 0x7F, 0x2C)
AMBER_100 = RGBColor(0xF1, 0xDD, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = E(1600)
prs.slide_height = E(940)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(0), E(0), E(1600), E(940))
bg.fill.solid()
bg.fill.fore_color.rgb = PAPER
bg.line.fill.background()
bg.shadow.inherit = False


def add_box(x, y, w, h, fill=None, line_color=NAVY_500, line_w=1.5,
            dashed=False, rounded=True, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, E(x), E(y), E(w), E(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
        if dashed:
            shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if radius is not None and rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.text_frame.margin_left = 0
    shp.text_frame.margin_right = 0
    shp.text_frame.margin_top = 0
    shp.text_frame.margin_bottom = 0
    return shp


def add_lines(shp, lines):
    """lines: list of (text, size, bold, color, align)"""
    tf = shp.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color


def add_text(x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    box.text_frame.vertical_anchor = anchor
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    add_lines(box, lines)
    return box


def add_arrow(points, color=NAVY_500, width=2.0, dashed=False, tail=True, head=False):
    xs = [p[0] for p in points]
    ys = [p[1] for p in points]
    builder = slide.shapes.build_freeform(E(xs[0]), E(ys[0]), scale=1)
    builder.add_line_segments([(E(px), E(py)) for px, py in points[1:]], close=False)
    shp = builder.convert_to_shape()
    shp.fill.background()
    shp.shadow.inherit = False
    shp.line.color.rgb = color
    shp.line.width = Pt(width)
    if dashed:
        shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    spPr = shp._element.spPr
    ln = spPr.find(qn('a:ln'))
    if tail:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if head:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return shp


def add_line(x1, y1, x2, y2, color=NAVY_300, width=1.0, dashed=False):
    ln = slide.shapes.add_connector(1, E(x1), E(y1), E(x2), E(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dashed:
        ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    ln.shadow.inherit = False
    return ln


CENTER = PP_ALIGN.CENTER
LEFT = PP_ALIGN.LEFT

# ---------------------------------------------------------------- header --
add_text(90, 8, 260, 30, [("AIGRE", 20, True, NAVY_900, LEFT)])
add_text(90, 40, 320, 20, [("AI GRIEVANCE RESOLUTION ENGINE", 8, True, AMBER_600, LEFT)])
add_text(880, 4, 660, 40, [("System Architecture", 24, True, NAVY_900, PP_ALIGN.RIGHT)])
add_text(880, 44, 660, 26, [("One Spring Boot application — intake through a human-approved decision",
                             11, False, NAVY_500, PP_ALIGN.RIGHT)])
add_line(60, 86, 1540, 86, NAVY_300, 0.75)

# --------------------------------------------------------- column A: UI ---
client = add_box(60, 180, 220, 270, fill=WHITE, line_color=NAVY_700, line_w=2)
add_lines(client, [
    ("Angular SPA", 14, True, NAVY_900, CENTER),
    ("Citizen Portal & Employee Dashboard", 10, False, NAVY_500, CENTER),
    ("submit · track · ask · manage", 9, False, NAVY_300, CENTER),
])
client.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

mailbox = add_box(60, 480, 220, 170, fill=WHITE, line_color=NAVY_500, line_w=2, dashed=True)
add_lines(mailbox, [
    ("Monitored Inbox", 14, True, NAVY_900, CENTER),
    ("Scheduled poll", 10, False, NAVY_500, CENTER),
    ("IMAP · unread only", 9, False, NAVY_300, CENTER),
])
mailbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# ------------------------------------------------------- backend panel ----
panel = add_box(320, 150, 860, 620, fill=None, line_color=NAVY_700, line_w=2.5, radius=0.03)
add_text(344, 158, 500, 20, [("AIGRE BACKEND", 12, True, NAVY_700, LEFT)])
add_text(344, 176, 700, 18, [("Spring Boot (WebFlux) — one application, four capabilities",
                              10, False, NAVY_500, LEFT)])
add_line(344, 200, 1156, 200, PAPER_LINE, 1)

intake = add_box(345, 232, 390, 220, fill=WHITE, line_color=NAVY_500, line_w=1.5)
add_lines(intake, [("Intake", 14, True, NAVY_900, LEFT),
                    ("Duplicate check · SLA due-date", 10, False, NAVY_500, LEFT)])
intake.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

workflow = add_box(765, 232, 380, 220, fill=AMBER_100, line_color=AMBER_600, line_w=2)
add_lines(workflow, [
    ("AI Classification", 14, True, NAVY_900, LEFT),
    ("& Human Review", 14, True, NAVY_900, LEFT),
    ("Pauses when confidence is low", 10, False, NAVY_700, LEFT),
])
workflow.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

ragchat = add_box(345, 472, 390, 220, fill=WHITE, line_color=NAVY_500, line_w=1.5)
add_lines(ragchat, [("Retrieval-Augmented Chat", 14, True, NAVY_900, LEFT),
                     ("Hybrid search + rerank, cited answers", 10, False, NAVY_500, LEFT)])
ragchat.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

dash = add_box(765, 472, 380, 220, fill=WHITE, line_color=NAVY_500, line_w=1.5)
add_lines(dash, [
    ("Dashboard Queries", 14, True, NAVY_900, LEFT),
    ("& MCP Tools", 14, True, NAVY_900, LEFT),
    ("Queues, Trends, status API", 10, False, NAVY_500, LEFT),
])
dash.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# indent the text boxes a bit from the box edge
for box in (intake, workflow, ragchat, dash):
    box.text_frame.margin_left = E(20)
    box.text_frame.margin_right = E(15)

# ------------------------------------------------------------- column C ---
data = add_box(1220, 180, 320, 270, fill=WHITE, line_color=NAVY_700, line_w=2)
add_lines(data, [
    ("PostgreSQL 16", 15, True, NAVY_900, CENTER),
    ("grievances · citizens · status", 10, False, NAVY_500, CENTER),
    ("+ pgvector", 13, True, NAVY_900, CENTER),
    ("policy corpus (vector search)", 10, False, NAVY_500, CENTER),
])
data.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

llm = add_box(1220, 480, 320, 220, fill=WHITE, line_color=NAVY_700, line_w=2)
add_text(1220, 500, 320, 24, [("LLM Providers", 13, True, NAVY_900, CENTER)])

ollama = add_box(1244, 536, 130, 48, fill=PAPER, line_color=NAVY_500, line_w=1.5, radius=0.15)
add_lines(ollama, [("Ollama", 10, True, NAVY_900, CENTER), ("local · default", 8, False, NAVY_500, CENTER)])
ollama.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

anthropic = add_box(1386, 536, 130, 48, fill=PAPER, line_color=NAVY_500, line_w=1.5, radius=0.15)
add_lines(anthropic, [("Anthropic", 10, True, NAVY_900, CENTER), ("Claude · swap", 8, False, NAVY_500, CENTER)])
anthropic.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

add_line(1374, 560, 1386, 560, NAVY_300, 1.25, dashed=True)
add_text(1220, 596, 320, 20, [("one config line switches providers", 9, False, NAVY_300, CENTER)])

# ------------------------------------------------------------- arrows -----
add_arrow([(280, 270), (345, 300)])
add_arrow([(280, 520), (302, 520), (302, 400), (345, 400)])
add_text(300, 452, 60, 16, [("poll", 9.5, False, NAVY_500, LEFT)])

add_arrow([(735, 232), (735, 221), (765, 221), (765, 232)])
add_text(715, 202, 90, 16, [("classify", 10, False, NAVY_500, CENTER)])

add_arrow([(945, 452), (945, 470)], color=AMBER_600, width=2.25)
add_text(955, 452, 110, 16, [("pending review", 11, True, AMBER_600, LEFT)])

add_arrow([(1000, 470), (1000, 452)], color=AMBER_600, width=2.25)
add_text(1010, 452, 90, 16, [("resumes", 11, True, AMBER_600, LEFT)])

add_arrow([(1180, 300), (1220, 270)])
add_text(1182, 278, 60, 16, [("commit", 10, False, NAVY_500, LEFT)])

add_arrow([(1180, 500), (1200, 500), (1200, 495), (1220, 495)])
add_text(1182, 504, 70, 16, [("classify", 10, False, NAVY_500, LEFT)])

# -------------------------------------------------------------- legend ----
leg1 = add_box(60, 880, 16, 16, fill=AMBER_100, line_color=AMBER_600, line_w=1.5, rounded=True, radius=0.2)
add_text(84, 878, 240, 20, [("AI pauses for a human decision", 12, False, NAVY_700, LEFT)])
add_line(330, 888, 366, 888, NAVY_300, 1.5, dashed=True)
add_text(376, 878, 160, 20, [("Swappable provider", 12, False, NAVY_700, LEFT)])
add_arrow([(540, 888), (576, 888)], color=NAVY_500, width=2)
add_text(586, 878, 140, 20, [("Data flow", 12, False, NAVY_700, LEFT)])

# ------------------------------------------------------- speaker notes ----
notes = slide.notes_slide
notes.notes_text_frame.text = (
    "Simplified from the full component diagram in docs/ARCHITECTURE.md. "
    "The \"commit\" and \"classify\" arrows into Postgres/LLM Providers represent "
    "every backend capability's traffic to them, not only the box each is drawn "
    "nearest to -- Dashboard's reads and Chat's retrieval/embedding calls travel "
    "the same two arrows. All shapes here are native and editable -- ungroup "
    "freely, recolor, or move anything."
)

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AIGRE-System-Architecture.pptx")
prs.save(out)
print("saved", out)
